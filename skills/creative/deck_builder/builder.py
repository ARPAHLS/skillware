"""Core PowerPoint (.pptx) builder, validator, and inspector for creative/deck_builder v0.2."""

from __future__ import annotations

import base64
import io
import json
import os
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import jsonschema
from PIL import Image, ImageOps
import pptx
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

try:
    from .placeholders import generate_placeholder
except ImportError:
    import sys

    sys.path.insert(0, os.path.dirname(__file__))
    from placeholders import generate_placeholder

_HERE = Path(__file__).resolve().parent
_TEMPLATES_DIR = _HERE / "templates"
_SCHEMA_PATH = _HERE / "schemas" / "deck_spec.schema.json"

TEMPLATES: Dict[str, Dict[str, Any]] = {
    "pitch_v1": {
        "template_id": "pitch_v1",
        "name": "Modern Pitch Deck",
        "description": "Vibrant modern aesthetic with bold typography and purple/indigo accents.",
        "aspect_ratio": "16:9",
        "default_accent": "#6E57E0",
        "default_heading_font": "Calibri",
        "default_body_font": "Calibri",
        "filename": "pitch_v1.pptx",
    },
    "corporate_v1": {
        "template_id": "corporate_v1",
        "name": "Executive Corporate",
        "description": "Structured executive presentation layout with navy and slate accents.",
        "aspect_ratio": "16:9",
        "default_accent": "#1E3A8A",
        "default_heading_font": "Calibri",
        "default_body_font": "Calibri",
        "filename": "corporate_v1.pptx",
    },
    "minimal_v1": {
        "template_id": "minimal_v1",
        "name": "Clean Minimalist",
        "description": "Monochrome high-contrast typography-forward layout for technical decks.",
        "aspect_ratio": "16:9",
        "default_accent": "#262626",
        "default_heading_font": "Arial",
        "default_body_font": "Arial",
        "filename": "minimal_v1.pptx",
    },
}

SUPPORTED_LAYOUT_TYPES = [
    "title",
    "section",
    "bullets",
    "two_column",
    "image",
    "image_caption",
    "quote",
    "table",
    "chart",
    "timeline",
    "metrics",
    "comparison",
    "blank",
]


def _load_schema() -> Dict[str, Any]:
    with open(_SCHEMA_PATH, "r", encoding="utf-8") as f:
        return json.load(f)


def _hex_to_rgb(hex_str: str) -> RGBColor:
    hex_clean = hex_str.lstrip("#")
    if len(hex_clean) == 3:
        hex_clean = "".join(c * 2 for c in hex_clean)
    if len(hex_clean) != 6:
        return RGBColor(110, 87, 224)
    r = int(hex_clean[0:2], 16)
    g = int(hex_clean[2:4], 16)
    b = int(hex_clean[4:6], 16)
    return RGBColor(r, g, b)


def list_templates() -> Dict[str, Any]:
    """Return bundled templates, descriptions, and aspect ratios."""
    template_list = []
    for tid, info in TEMPLATES.items():
        template_list.append(
            {
                "template_id": tid,
                "name": info["name"],
                "description": info["description"],
                "aspect_ratio": info["aspect_ratio"],
                "default_accent": info["default_accent"],
                "supported_layouts": SUPPORTED_LAYOUT_TYPES,
            }
        )
    return {
        "success": True,
        "action": "list_templates",
        "templates": template_list,
        "count": len(template_list),
    }


def _validate_output_path(output_path: str) -> str:
    if not output_path or not output_path.strip():
        raise ValueError("output_path must be a non-empty string.")

    raw_parts = output_path.replace("\\", "/").split("/")
    if ".." in raw_parts:
        raise ValueError(
            "output_path contains prohibited path traversal sequences ('..')."
        )

    norm = os.path.normpath(output_path)
    if "\x00" in norm:
        raise ValueError("output_path contains invalid characters.")

    if not norm.lower().endswith(".pptx"):
        raise ValueError("output_path must have a .pptx extension.")

    parent = os.path.dirname(os.path.abspath(norm))
    os.makedirs(parent, exist_ok=True)
    return os.path.abspath(norm)


def _load_raw_image_bytes(
    img_obj: Dict[str, Any], fallback_kind: str = "hero"
) -> Optional[bytes]:
    if not img_obj or not isinstance(img_obj, dict):
        return None
    if img_obj.get("path"):
        p = img_obj["path"]
        if os.path.exists(p):
            with open(p, "rb") as f:
                return f.read()
    elif img_obj.get("base64"):
        try:
            return base64.b64decode(img_obj["base64"])
        except Exception:
            return None
    elif img_obj.get("placeholder_id") or img_obj.get("placeholder_prompt"):
        p_stream = generate_placeholder(
            kind=img_obj.get("placeholder_id", fallback_kind),
            prompt=img_obj.get("placeholder_prompt"),
        )
        return p_stream.getvalue()
    return None


def _normalize_and_fit_image(
    raw_bytes: bytes,
    box_width: float,
    box_height: float,
    fit: str = "contain",
) -> Tuple[io.BytesIO, float, float, float, float]:
    """
    Normalizes image with Pillow and computes fitting offsets.
    Returns: (img_stream, offset_x, offset_y, draw_w, draw_h)
    """
    fit = (fit or "contain").lower().strip()
    with Image.open(io.BytesIO(raw_bytes)) as pil_img:
        pil_img = ImageOps.exif_transpose(pil_img) or pil_img
        if pil_img.mode in ("CMYK", "P"):
            pil_img = pil_img.convert(
                "RGBA"
                if "transparency" in pil_img.info or pil_img.mode == "P"
                else "RGB"
            )

        w, h = pil_img.size
        img_aspect = float(w) / float(h) if h > 0 else 1.0
        box_aspect = float(box_width) / float(box_height) if box_height > 0 else 1.0

        if fit in ("cover", "crop_center"):
            if img_aspect > box_aspect:
                new_w = int(h * box_aspect)
                ox = (w - new_w) // 2
                pil_img = pil_img.crop((ox, 0, ox + new_w, h))
            elif img_aspect < box_aspect:
                new_h = int(w / box_aspect)
                oy = (h - new_h) // 2
                pil_img = pil_img.crop((0, oy, w, oy + new_h))
            draw_w = box_width
            draw_h = box_height
            offset_x = 0.0
            offset_y = 0.0
        elif fit == "contain":
            if img_aspect > box_aspect:
                draw_w = box_width
                draw_h = box_width / img_aspect
                offset_x = 0.0
                offset_y = (box_height - draw_h) / 2.0
            else:
                draw_h = box_height
                draw_w = box_height * img_aspect
                offset_y = 0.0
                offset_x = (box_width - draw_w) / 2.0
        else:  # stretch or native
            draw_w = box_width
            draw_h = box_height
            offset_x = 0.0
            offset_y = 0.0

        buf = io.BytesIO()
        pil_img.save(buf, format="PNG")
        buf.seek(0)
        return buf, offset_x, offset_y, draw_w, draw_h


def validate_spec(deck_spec: Dict[str, Any], strict: bool = False) -> Dict[str, Any]:
    """Validate deck specification against schema and business rules."""
    schema = _load_schema()
    validator = jsonschema.Draft202012Validator(schema)
    schema_errors = list(validator.iter_errors(deck_spec))

    errors: List[Dict[str, Any]] = []
    warnings: List[Dict[str, Any]] = []

    for err in schema_errors:
        path = ".".join(str(p) for p in err.path) or "root"
        errors.append(
            {
                "code": "SCHEMA_VALIDATION_ERROR",
                "slide_index": -1,
                "message": f"[{path}] {err.message}",
            }
        )

    if errors:
        return {
            "success": False,
            "action": "validate_spec",
            "valid": False,
            "template_id": deck_spec.get("template_id", "unknown"),
            "slide_count": len(deck_spec.get("slides", [])),
            "warnings": [],
            "errors": errors,
            "error_code": "INVALID_SPEC",
        }

    template_id = deck_spec.get("template_id", "pitch_v1")
    if template_id not in TEMPLATES:
        warnings.append(
            {
                "code": "UNKNOWN_TEMPLATE_FALLBACK",
                "slide_index": -1,
                "message": f"Template '{template_id}' is not bundled. Falling back to 'pitch_v1'.",
            }
        )

    slides = deck_spec.get("slides", [])
    for idx, slide in enumerate(slides):
        stype = slide.get("type")
        if stype not in SUPPORTED_LAYOUT_TYPES:
            errors.append(
                {
                    "code": "UNSUPPORTED_LAYOUT_TYPE",
                    "slide_index": idx,
                    "message": f"Layout type '{stype}' is not supported.",
                }
            )
            continue

        # Soft cap check on bullets length
        if stype == "bullets":
            bullets = slide.get("bullets", [])
            for b_idx, bullet in enumerate(bullets):
                if len(str(bullet)) > 120:
                    warnings.append(
                        {
                            "code": "BULLET_TRUNCATED",
                            "slide_index": idx,
                            "bullet_index": b_idx,
                            "character_count": len(str(bullet)),
                            "message": (
                                f"Bullet {b_idx + 1} exceeds 120 characters "
                                f"({len(str(bullet))} chars). Text may wrap excessively."
                            ),
                        }
                    )

        # Asset verification for images
        if stype in {"image", "image_caption", "title"}:
            img_obj = slide.get("image")
            if img_obj and isinstance(img_obj, dict):
                img_path = img_obj.get("path")
                b64_data = img_obj.get("base64")
                p_id = img_obj.get("placeholder_id")
                p_prompt = img_obj.get("placeholder_prompt")
                if img_path:
                    if not os.path.exists(img_path):
                        warnings.append(
                            {
                                "code": "ASSET_NOT_FOUND",
                                "slide_index": idx,
                                "message": f"Image file not found at path: {img_path}",
                            }
                        )
                elif b64_data:
                    try:
                        raw_bytes = base64.b64decode(b64_data)
                        with Image.open(io.BytesIO(raw_bytes)) as pil_img:
                            pil_img.verify()
                    except Exception as exc:
                        warnings.append(
                            {
                                "code": "ASSET_INVALID",
                                "slide_index": idx,
                                "message": f"Invalid base64 image data: {exc}",
                            }
                        )
                elif not p_id and not p_prompt:
                    warnings.append(
                        {
                            "code": "ASSET_SPEC_EMPTY",
                            "slide_index": idx,
                            "message": "Image object specified but lacks 'path', 'base64', or 'placeholder_id'.",
                        }
                    )

        # Chart verification
        if stype == "chart":
            chart_obj = slide.get("chart") or {}
            cats = chart_obj.get("categories") or []
            series_list = chart_obj.get("series") or []
            for s in series_list:
                vals = s.get("values") or []
                if len(vals) != len(cats):
                    errors.append(
                        {
                            "code": "CHART_DIMENSION_MISMATCH",
                            "slide_index": idx,
                            "message": (
                                f"Series '{s.get('name')}' length ({len(vals)}) "
                                f"does not match categories count ({len(cats)})."
                            ),
                        }
                    )

        # Table verification
        if stype == "table":
            cols = slide.get("columns") or []
            rows = slide.get("rows") or []
            for r_idx, row in enumerate(rows):
                if len(row) != len(cols):
                    warnings.append(
                        {
                            "code": "TABLE_DIMENSION_MISMATCH",
                            "slide_index": idx,
                            "message": (
                                f"Row {r_idx + 1} length ({len(row)}) does not match "
                                f"column headers count ({len(cols)})."
                            ),
                        }
                    )

        # Timeline verification
        if stype == "timeline":
            items = slide.get("items") or []
            if not items:
                errors.append(
                    {
                        "code": "MISSING_REQUIRED_FIELD",
                        "slide_index": idx,
                        "message": "Timeline slide requires at least one milestone in 'items'.",
                    }
                )
            for it_idx, item in enumerate(items):
                if not item.get("date") or not item.get("title"):
                    errors.append(
                        {
                            "code": "INVALID_TIMELINE_ITEM",
                            "slide_index": idx,
                            "message": f"Timeline milestone {it_idx + 1} must include 'date' and 'title'.",
                        }
                    )

        # Metrics verification
        if stype == "metrics":
            m_list = slide.get("metrics") or []
            if not m_list:
                errors.append(
                    {
                        "code": "MISSING_REQUIRED_FIELD",
                        "slide_index": idx,
                        "message": "Metrics slide requires at least one metric in 'metrics'.",
                    }
                )
            for m_idx, m in enumerate(m_list):
                if not m.get("value") or not m.get("label"):
                    errors.append(
                        {
                            "code": "INVALID_METRIC_ITEM",
                            "slide_index": idx,
                            "message": f"Metric {m_idx + 1} must include 'value' and 'label'.",
                        }
                    )

    if strict and warnings:
        for w in warnings:
            errors.append(
                {
                    "code": f"STRICT_{w['code']}",
                    "slide_index": w.get("slide_index", -1),
                    "message": f"Strict mode validation failure: {w['message']}",
                }
            )

    is_valid = len(errors) == 0
    return {
        "success": is_valid,
        "action": "validate_spec",
        "valid": is_valid,
        "template_id": template_id,
        "slide_count": len(slides),
        "warnings": warnings,
        "errors": errors,
        "error_code": (
            None if is_valid else (errors[0]["code"] if errors else "VALIDATION_FAILED")
        ),
    }


def render_deck(
    deck_spec: Dict[str, Any],
    output_path: str,
    template_id: Optional[str] = None,
    theme: Optional[Dict[str, Any]] = None,
    strict: bool = False,
) -> Dict[str, Any]:
    """Render presentation from deck specification to target .pptx file."""
    try:
        safe_output_path = _validate_output_path(output_path)
    except ValueError as exc:
        return {
            "success": False,
            "action": "render",
            "output_path": output_path,
            "template_id": template_id or "unknown",
            "slide_count": 0,
            "file_size_bytes": 0,
            "slides": [],
            "warnings": [],
            "errors": [
                {"code": "OUTPUT_PATH_UNSAFE", "slide_index": -1, "message": str(exc)}
            ],
            "error_code": "OUTPUT_PATH_UNSAFE",
        }

    val_res = validate_spec(deck_spec, strict=strict)
    if not val_res["valid"]:
        return {
            "success": False,
            "action": "render",
            "output_path": safe_output_path,
            "template_id": val_res["template_id"],
            "slide_count": val_res["slide_count"],
            "file_size_bytes": 0,
            "slides": [],
            "warnings": val_res["warnings"],
            "errors": val_res["errors"],
            "error_code": val_res["error_code"] or "VALIDATION_FAILED",
        }

    effective_template_id = template_id or deck_spec.get("template_id") or "pitch_v1"
    if effective_template_id not in TEMPLATES:
        effective_template_id = "pitch_v1"

    template_info = TEMPLATES[effective_template_id]
    template_path = _TEMPLATES_DIR / template_info["filename"]

    if not template_path.exists():
        prs = pptx.Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = pptx.Presentation(str(template_path))

    # Determine theme tokens
    theme_spec = theme or deck_spec.get("theme") or {}
    accent_hex = theme_spec.get("accent_color") or template_info["default_accent"]
    font_heading = (
        theme_spec.get("font_heading") or template_info["default_heading_font"]
    )
    font_body = theme_spec.get("font_body") or template_info["default_body_font"]
    accent_rgb = _hex_to_rgb(accent_hex)

    # Document governance metadata (prefer metadata, fallback to root)
    metadata = deck_spec.get("metadata") or {}
    classification = metadata.get("classification") or deck_spec.get("classification")
    legal_footer = metadata.get("legal_footer") or deck_spec.get("legal_footer")

    slides = deck_spec.get("slides", [])
    rendered_slides_summary: List[Dict[str, Any]] = []

    for s_idx, slide_data in enumerate(slides):
        stype = slide_data.get("type")
        title_text = slide_data.get("title", "")
        slide = None

        # 1. Title Slide
        if stype == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = font_heading
                    p.font.color.rgb = accent_rgb

            subtitle_text = slide_data.get("subtitle", "")
            if subtitle_text and len(slide.placeholders) > 1:
                sub_ph = slide.placeholders[1]
                sub_ph.text = subtitle_text
                for p in sub_ph.text_frame.paragraphs:
                    p.font.name = font_body

            raw_bytes = _load_raw_image_bytes(
                slide_data.get("image"), fallback_kind="logo"
            )
            if raw_bytes:
                fit_mode = slide_data.get("image", {}).get("fit", "contain")
                stream, ox, oy, dw, dh = _normalize_and_fit_image(
                    raw_bytes, 2.0, 2.0, fit=fit_mode
                )
                try:
                    slide.shapes.add_picture(
                        stream,
                        Inches(1.2 + ox),
                        Inches(1.0 + oy),
                        width=Inches(dw),
                        height=Inches(dh),
                    )
                except Exception:
                    pass

        # 2. Section Header
        elif stype == "section":
            slide = prs.slides.add_slide(prs.slide_layouts[2])
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = font_heading
                    p.font.color.rgb = accent_rgb

            subtitle_text = slide_data.get("subtitle", "")
            if subtitle_text and len(slide.placeholders) > 1:
                sub_ph = slide.placeholders[1]
                sub_ph.text = subtitle_text
                for p in sub_ph.text_frame.paragraphs:
                    p.font.name = font_body

        # 3. Bullets Slide
        elif stype == "bullets":
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = font_heading
                    p.font.color.rgb = accent_rgb

            bullets = slide_data.get("bullets", [])
            if len(slide.placeholders) > 1:
                body_ph = slide.placeholders[1]
                tf = body_ph.text_frame
                tf.word_wrap = True
                for b_idx, bullet in enumerate(bullets):
                    if b_idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = str(bullet)
                    p.font.name = font_body
                    p.level = 0

        # 4. Two Column Slide
        elif stype == "two_column":
            slide = prs.slides.add_slide(prs.slide_layouts[3])
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = font_heading
                    p.font.color.rgb = accent_rgb

            left_content = slide_data.get("left", [])
            if len(slide.placeholders) > 1:
                left_ph = slide.placeholders[1]
                tf_l = left_ph.text_frame
                tf_l.word_wrap = True
                if isinstance(left_content, list):
                    for idx_item, item in enumerate(left_content):
                        p = (
                            tf_l.paragraphs[0]
                            if idx_item == 0
                            else tf_l.add_paragraph()
                        )
                        p.text = str(item)
                        p.font.name = font_body
                else:
                    tf_l.paragraphs[0].text = str(left_content)
                    tf_l.paragraphs[0].font.name = font_body

            right_content = slide_data.get("right", [])
            if len(slide.placeholders) > 2:
                right_ph = slide.placeholders[2]
                tf_r = right_ph.text_frame
                tf_r.word_wrap = True
                if isinstance(right_content, list):
                    for idx_item, item in enumerate(right_content):
                        p = (
                            tf_r.paragraphs[0]
                            if idx_item == 0
                            else tf_r.add_paragraph()
                        )
                        p.text = str(item)
                        p.font.name = font_body
                else:
                    tf_r.paragraphs[0].text = str(right_content)
                    tf_r.paragraphs[0].font.name = font_body

        # 5. Image Slide
        elif stype == "image":
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if slide.shapes.title and title_text:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = font_heading
                    p.font.color.rgb = accent_rgb

            raw_bytes = _load_raw_image_bytes(
                slide_data.get("image"), fallback_kind="hero"
            )
            if raw_bytes:
                fit_mode = slide_data.get("image", {}).get("fit", "contain")
                box_w = 9.333
                box_h = 4.2 if title_text else 5.0
                stream, ox, oy, dw, dh = _normalize_and_fit_image(
                    raw_bytes, box_w, box_h, fit=fit_mode
                )
                top_offset = 1.8 if title_text else 1.2
                try:
                    slide.shapes.add_picture(
                        stream,
                        Inches(2.0 + ox),
                        Inches(top_offset + oy),
                        width=Inches(dw),
                        height=Inches(dh),
                    )
                except Exception:
                    pass

            caption = slide_data.get("caption")
            if caption:
                tb = slide.shapes.add_textbox(
                    Inches(2.0), Inches(6.2), Inches(9.333), Inches(0.8)
                )
                p = tb.text_frame.paragraphs[0]
                p.text = caption
                p.font.name = font_body
                p.font.italic = True

        # 6. Image with Caption Body
        elif stype == "image_caption":
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = font_heading
                    p.font.color.rgb = accent_rgb

            raw_bytes = _load_raw_image_bytes(
                slide_data.get("image"), fallback_kind="hero"
            )
            if raw_bytes:
                fit_mode = slide_data.get("image", {}).get("fit", "contain")
                stream, ox, oy, dw, dh = _normalize_and_fit_image(
                    raw_bytes, 5.5, 4.5, fit=fit_mode
                )
                try:
                    slide.shapes.add_picture(
                        stream,
                        Inches(1.0 + ox),
                        Inches(1.8 + oy),
                        width=Inches(dw),
                        height=Inches(dh),
                    )
                except Exception:
                    pass

            body_text = slide_data.get("body", "")
            if body_text:
                tb = slide.shapes.add_textbox(
                    Inches(7.0), Inches(1.8), Inches(5.333), Inches(4.5)
                )
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = body_text
                p.font.name = font_body

        # 7. Quote Slide
        elif stype == "quote":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tb = slide.shapes.add_textbox(
                Inches(1.8), Inches(2.0), Inches(9.7), Inches(3.5)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            quote_text = slide_data.get("quote", "")
            p_q = tf.paragraphs[0]
            p_q.text = f'"{quote_text}"'
            p_q.font.name = font_heading
            p_q.font.size = Pt(26)
            p_q.font.bold = True
            p_q.font.color.rgb = accent_rgb
            attrib_text = slide_data.get("attribution", "")
            if attrib_text:
                p_a = tf.add_paragraph()
                p_a.text = f"— {attrib_text}"
                p_a.font.name = font_body
                p_a.font.size = Pt(18)
                p_a.font.color.rgb = RGBColor(100, 116, 139)

        # 8. Table Slide
        elif stype == "table":
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = font_heading
                    p.font.color.rgb = accent_rgb

            cols = slide_data.get("columns", [])
            rows = slide_data.get("rows", [])
            num_rows = len(rows) + 1
            num_cols = len(cols)
            if num_cols > 0:
                table_shape = slide.shapes.add_table(
                    num_rows,
                    num_cols,
                    Inches(1.2),
                    Inches(1.8),
                    Inches(10.933),
                    Inches(0.6 * num_rows),
                )
                tbl = table_shape.table
                for c_idx, col_name in enumerate(cols):
                    cell = tbl.cell(0, c_idx)
                    cell.text = str(col_name)
                    for p in cell.text_frame.paragraphs:
                        p.font.name = font_heading
                        p.font.bold = True
                        p.font.color.rgb = accent_rgb
                for r_idx, row in enumerate(rows):
                    for c_idx, val in enumerate(row):
                        if c_idx < num_cols:
                            cell = tbl.cell(r_idx + 1, c_idx)
                            cell.text = str(val)
                            for p in cell.text_frame.paragraphs:
                                p.font.name = font_body

        # 9. Chart Slide
        elif stype == "chart":
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = font_heading
                    p.font.color.rgb = accent_rgb

            chart_obj = slide_data.get("chart", {})
            kind = chart_obj.get("kind", "bar").lower()
            categories = chart_obj.get("categories", [])
            series_list = chart_obj.get("series", [])

            chart_type_map = {
                "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "line": XL_CHART_TYPE.LINE_MARKERS,
                "pie": XL_CHART_TYPE.PIE,
            }
            xl_type = chart_type_map.get(kind, XL_CHART_TYPE.COLUMN_CLUSTERED)

            chart_data = CategoryChartData()
            chart_data.categories = [str(c) for c in categories]
            for s in series_list:
                chart_data.add_series(str(s.get("name", "")), s.get("values", []))

            slide.shapes.add_chart(
                xl_type,
                Inches(1.5),
                Inches(1.8),
                Inches(10.333),
                Inches(4.8),
                chart_data,
            )

        # 10. Timeline Slide (NEW v0.2)
        elif stype == "timeline":
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = font_heading
                    p.font.color.rgb = accent_rgb

            items = slide_data.get("items", [])
            n_items = len(items)
            if n_items > 0:
                avail_w = 10.933
                step_w = avail_w / n_items
                line_y = 3.3

                # Horizontal bar
                connector = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1.2),
                    Inches(line_y),
                    Inches(avail_w),
                    Inches(0.04),
                )
                connector.fill.solid()
                connector.fill.fore_color.rgb = RGBColor(226, 232, 240)
                connector.line.fill.background()

                for i, it in enumerate(items):
                    item_x = 1.2 + i * step_w
                    card_w = step_w - 0.2

                    # Status dot/pill
                    dot = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL,
                        Inches(item_x + card_w / 2 - 0.12),
                        Inches(line_y - 0.1),
                        Inches(0.24),
                        Inches(0.24),
                    )
                    dot.fill.solid()
                    status = it.get("status", "planned").lower()
                    if status == "completed":
                        dot.fill.fore_color.rgb = RGBColor(16, 185, 129)
                    elif status == "in_progress":
                        dot.fill.fore_color.rgb = accent_rgb
                    else:
                        dot.fill.fore_color.rgb = RGBColor(148, 163, 184)
                    dot.line.fill.background()

                    # Date text (above)
                    tb_date = slide.shapes.add_textbox(
                        Inches(item_x), Inches(2.2), Inches(card_w), Inches(0.8)
                    )
                    p_d = tb_date.text_frame.paragraphs[0]
                    p_d.text = str(it.get("date", ""))
                    p_d.font.name = font_heading
                    p_d.font.bold = True
                    p_d.font.size = Pt(14)
                    p_d.font.color.rgb = accent_rgb

                    # Title & description (below)
                    tb_desc = slide.shapes.add_textbox(
                        Inches(item_x), Inches(3.7), Inches(card_w), Inches(2.5)
                    )
                    tf_d = tb_desc.text_frame
                    tf_d.word_wrap = True
                    p_t = tf_d.paragraphs[0]
                    p_t.text = str(it.get("title", ""))
                    p_t.font.name = font_heading
                    p_t.font.bold = True
                    p_t.font.size = Pt(14)
                    if it.get("description"):
                        p_sub = tf_d.add_paragraph()
                        p_sub.text = str(it.get("description"))
                        p_sub.font.name = font_body
                        p_sub.font.size = Pt(11)
                        p_sub.font.color.rgb = RGBColor(100, 116, 139)

        # 11. Metrics Slide (NEW v0.2)
        elif stype == "metrics":
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = font_heading
                    p.font.color.rgb = accent_rgb

            metrics_list = slide_data.get("metrics", [])
            n_metrics = min(len(metrics_list), 4)
            if n_metrics > 0:
                avail_w = 11.333
                gap = 0.35
                card_w = (avail_w - (n_metrics - 1) * gap) / n_metrics
                for i, m in enumerate(metrics_list[:n_metrics]):
                    cx = 1.0 + i * (card_w + gap)

                    # Card frame
                    card = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(cx),
                        Inches(2.2),
                        Inches(card_w),
                        Inches(4.2),
                    )
                    card.fill.solid()
                    card.fill.fore_color.rgb = RGBColor(248, 250, 252)
                    card.line.color.rgb = RGBColor(226, 232, 240)
                    card.line.width = Pt(1)

                    # Content
                    tb = slide.shapes.add_textbox(
                        Inches(cx + 0.15),
                        Inches(2.5),
                        Inches(card_w - 0.3),
                        Inches(3.6),
                    )
                    tf = tb.text_frame
                    tf.word_wrap = True

                    # Big number
                    p_val = tf.paragraphs[0]
                    val_str = str(m.get("value", ""))
                    if m.get("unit"):
                        val_str = f"{val_str} {m.get('unit')}"
                    p_val.text = val_str
                    p_val.font.name = font_heading
                    p_val.font.size = Pt(36)
                    p_val.font.bold = True
                    p_val.font.color.rgb = accent_rgb

                    # Trend indicator
                    delta = m.get("delta")
                    trend = m.get("trend", "neutral")
                    if delta:
                        p_t = tf.add_paragraph()
                        symbol = (
                            "↑ "
                            if trend == "up"
                            else ("↓ " if trend == "down" else "→ ")
                        )
                        p_t.text = f"{symbol}{delta}"
                        p_t.font.name = font_heading
                        p_t.font.size = Pt(13)
                        p_t.font.bold = True
                        if trend == "up":
                            p_t.font.color.rgb = RGBColor(16, 185, 129)
                        elif trend == "down":
                            p_t.font.color.rgb = RGBColor(239, 68, 68)
                        else:
                            p_t.font.color.rgb = RGBColor(100, 116, 139)

                    # Label
                    p_lbl = tf.add_paragraph()
                    p_lbl.text = str(m.get("label", ""))
                    p_lbl.font.name = font_body
                    p_lbl.font.size = Pt(14)
                    p_lbl.font.bold = True
                    p_lbl.font.color.rgb = RGBColor(71, 85, 105)

        # 12. Comparison Slide (NEW v0.2)
        elif stype == "comparison":
            slide = prs.slides.add_slide(prs.slide_layouts[3])
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                for p in slide.shapes.title.text_frame.paragraphs:
                    p.font.name = font_heading
                    p.font.color.rgb = accent_rgb

            def _populate_column(tf, content):
                tf.word_wrap = True
                if isinstance(content, dict):
                    col_title = content.get("title")
                    items = content.get("items", [])
                    started = False
                    if col_title:
                        p = tf.paragraphs[0]
                        p.text = str(col_title)
                        p.font.name = font_heading
                        p.font.bold = True
                        p.font.size = Pt(16)
                        p.font.color.rgb = accent_rgb
                        started = True
                    for item in items:
                        p = tf.paragraphs[0] if not started else tf.add_paragraph()
                        p.text = str(item)
                        p.font.name = font_body
                        started = True
                elif isinstance(content, list):
                    for idx_item, item in enumerate(content):
                        p = tf.paragraphs[0] if idx_item == 0 else tf.add_paragraph()
                        p.text = str(item)
                        p.font.name = font_body
                else:
                    tf.paragraphs[0].text = str(content)
                    tf.paragraphs[0].font.name = font_body

            left_content = slide_data.get("left", [])
            if len(slide.placeholders) > 1:
                _populate_column(slide.placeholders[1].text_frame, left_content)

            right_content = slide_data.get("right", [])
            if len(slide.placeholders) > 2:
                _populate_column(slide.placeholders[2].text_frame, right_content)

        # 13. Blank
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Attach speaker notes if present
        speaker_notes = slide_data.get("speaker_notes")
        if speaker_notes and slide:
            slide.notes_slide.notes_text_frame.text = speaker_notes

        # Governance classification ribbon & legal footer
        if slide:
            if classification:
                tb_c = slide.shapes.add_textbox(
                    Inches(4.5), Inches(0.1), Inches(4.333), Inches(0.3)
                )
                p_c = tb_c.text_frame.paragraphs[0]
                p_c.text = f"[{classification.upper()}]"
                p_c.font.name = font_heading
                p_c.font.size = Pt(8)
                p_c.font.bold = True
                p_c.font.color.rgb = RGBColor(148, 163, 184)
            if legal_footer:
                tb_f = slide.shapes.add_textbox(
                    Inches(0.8), Inches(7.1), Inches(11.7), Inches(0.3)
                )
                p_f = tb_f.text_frame.paragraphs[0]
                p_f.text = str(legal_footer)
                p_f.font.name = font_body
                p_f.font.size = Pt(8)
                p_f.font.color.rgb = RGBColor(148, 163, 184)

        rendered_slides_summary.append(
            {
                "index": s_idx,
                "type": stype,
                "title": title_text,
            }
        )

    prs.save(safe_output_path)
    file_size = os.path.getsize(safe_output_path)

    return {
        "success": True,
        "action": "render",
        "output_path": safe_output_path,
        "template_id": effective_template_id,
        "slide_count": len(slides),
        "file_size_bytes": file_size,
        "slides": rendered_slides_summary,
        "warnings": val_res["warnings"],
        "error_code": None,
    }


def inspect_deck(input_path: str) -> Dict[str, Any]:
    """Inspect an existing .pptx presentation and return its slide manifest."""
    if not input_path or not os.path.exists(input_path):
        return {
            "success": False,
            "action": "inspect",
            "slide_count": 0,
            "slides": [],
            "error_code": "INSPECT_FAILED",
            "errors": [
                {
                    "code": "FILE_NOT_FOUND",
                    "slide_index": -1,
                    "message": f"File not found at input_path: {input_path}",
                }
            ],
        }

    try:
        prs = pptx.Presentation(input_path)
        slides_manifest: List[Dict[str, Any]] = []
        for idx, s in enumerate(prs.slides):
            has_notes = False
            try:
                has_notes = bool(
                    s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip()
                )
            except Exception:
                pass
            title_text = (
                s.shapes.title.text
                if (s.shapes.title and s.shapes.title.text)
                else None
            )
            layout_name = getattr(s.slide_layout, "name", "Custom")
            slides_manifest.append(
                {
                    "index": idx,
                    "layout_name": layout_name,
                    "title": title_text,
                    "has_notes": has_notes,
                    "shape_count": len(s.shapes),
                }
            )

        return {
            "success": True,
            "action": "inspect",
            "slide_count": len(slides_manifest),
            "slides": slides_manifest,
            "error_code": None,
        }
    except Exception as exc:
        return {
            "success": False,
            "action": "inspect",
            "slide_count": 0,
            "slides": [],
            "error_code": "INSPECT_FAILED",
            "errors": [
                {
                    "code": "CORRUPT_OR_UNREADABLE",
                    "slide_index": -1,
                    "message": f"Failed to inspect .pptx file: {exc}",
                }
            ],
        }
