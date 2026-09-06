"""Unit and bundle tests for creative/deck_builder."""

import base64
import io
import os
from PIL import Image
import pytest

from skills.creative.deck_builder.skill import DeckBuilderSkill


@pytest.fixture
def skill():
    return DeckBuilderSkill()


@pytest.fixture
def sample_base64_png():
    img = Image.new("RGB", (100, 100), color=(110, 87, 224))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return base64.b64encode(buf.getvalue()).decode("utf-8")


def test_manifest_loads_and_declares_requirements(skill):
    manifest = skill.manifest
    assert manifest["name"] == "creative/deck_builder"
    assert manifest["version"] == "0.2.0"
    assert manifest["category"] == "creative"
    assert "python-pptx>=1.0.0" in manifest["requirements"]
    assert "pillow" in manifest["requirements"]
    assert manifest["issuer"]["org"] == "ARPAHLS"


def test_list_templates_returns_bundled_templates(skill):
    res = skill.execute({"action": "list_templates"})
    assert res["success"] is True
    assert len(res["templates"]) >= 3
    tids = [t["template_id"] for t in res["templates"]]
    assert "pitch_v1" in tids
    assert "corporate_v1" in tids
    assert "minimal_v1" in tids


def test_validate_spec_valid_payload(skill):
    spec = {
        "title": "Clean Pitch",
        "template_id": "pitch_v1",
        "slides": [
            {"type": "title", "title": "Overview", "subtitle": "A short summary"},
            {
                "type": "bullets",
                "title": "Key Points",
                "bullets": ["First point", "Second point"],
            },
        ],
    }
    res = skill.execute({"action": "validate_spec", "deck_spec": spec})
    assert res["success"] is True
    assert res["valid"] is True
    assert res["slide_count"] == 2
    assert len(res["errors"]) == 0


def test_validate_spec_invalid_schema_rejects(skill):
    bad_spec = {"title": 12345}  # missing slides and invalid title type
    res = skill.execute({"action": "validate_spec", "deck_spec": bad_spec})
    assert res["success"] is False
    assert res["valid"] is False
    assert res["error_code"] == "INVALID_SPEC"
    assert len(res["errors"]) > 0


def test_validate_spec_bullet_truncated_warning(skill):
    spec = {
        "title": "Lengthy Bullet Deck",
        "slides": [
            {
                "type": "bullets",
                "title": "Long Bullet Slide",
                "bullets": [
                    "Short bullet",
                    "A" * 130,  # exceeds 120 chars
                ],
            }
        ],
    }
    res = skill.execute({"action": "validate_spec", "deck_spec": spec})
    assert res["valid"] is True  # non-fatal by default
    assert any(w["code"] == "BULLET_TRUNCATED" for w in res["warnings"])


def test_validate_spec_strict_mode_fails_on_warning(skill):
    spec = {
        "title": "Strict Mode Test",
        "slides": [
            {
                "type": "bullets",
                "title": "Long Bullet Slide",
                "bullets": ["A" * 135],
            }
        ],
    }
    res = skill.execute({"action": "validate_spec", "deck_spec": spec, "strict": True})
    assert res["valid"] is False
    assert res["success"] is False
    assert any("STRICT_" in e["code"] for e in res["errors"])


def test_validate_spec_chart_dimension_mismatch(skill):
    spec = {
        "title": "Mismatched Chart",
        "slides": [
            {
                "type": "chart",
                "title": "Growth",
                "chart": {
                    "kind": "bar",
                    "categories": ["Q1", "Q2", "Q3"],
                    "series": [
                        {"name": "Users", "values": [10, 20]}
                    ],  # only 2 values for 3 categories
                },
            }
        ],
    }
    res = skill.execute({"action": "validate_spec", "deck_spec": spec})
    assert res["valid"] is False
    assert any(e["code"] == "CHART_DIMENSION_MISMATCH" for e in res["errors"])


def test_validate_spec_asset_not_found(skill):
    spec = {
        "title": "Missing Image Deck",
        "slides": [
            {
                "type": "image",
                "title": "Photo",
                "image": {"path": "/nonexistent/path/to/missing_file.png"},
            }
        ],
    }
    res = skill.execute({"action": "validate_spec", "deck_spec": spec})
    assert any(w["code"] == "ASSET_NOT_FOUND" for w in res["warnings"])


def test_validate_spec_invalid_base64_asset(skill):
    spec = {
        "title": "Corrupt Base64 Deck",
        "slides": [
            {
                "type": "image",
                "title": "Corrupt",
                "image": {"base64": "!!!not_valid_base64$$$"},
            }
        ],
    }
    res = skill.execute({"action": "validate_spec", "deck_spec": spec})
    assert any(w["code"] == "ASSET_INVALID" for w in res["warnings"])


def test_render_all_ten_slide_types(skill, tmp_path, sample_base64_png):
    img_file = tmp_path / "test_logo.png"
    Image.new("RGB", (80, 80), color=(20, 150, 80)).save(img_file)

    out_file = tmp_path / "all_types.pptx"
    spec = {
        "title": "Full Feature Deck",
        "template_id": "corporate_v1",
        "theme": {
            "accent_color": "#1E3A8A",
            "font_heading": "Calibri",
            "font_body": "Calibri",
        },
        "slides": [
            {
                "type": "title",
                "title": "Corporate Briefing",
                "subtitle": "Executive summary",
                "image": {"path": str(img_file)},
            },
            {
                "type": "section",
                "title": "Part 1: Operational Review",
                "subtitle": "Key metrics and progress",
            },
            {
                "type": "bullets",
                "title": "Strategic Objectives",
                "bullets": ["Expand market presence", "Optimize unit economics"],
            },
            {
                "type": "two_column",
                "title": "Comparison",
                "left": ["Legacy Approach", "Manual checks"],
                "right": ["Skillware", "Autonomous validation"],
            },
            {
                "type": "image",
                "title": "Visual Dashboard",
                "image": {"base64": sample_base64_png},
                "caption": "Figure 1: Pipeline efficiency",
            },
            {
                "type": "image_caption",
                "title": "Architecture",
                "image": {"path": str(img_file)},
                "body": "Microservices communicate via gRPC with strict contracts.",
            },
            {
                "type": "quote",
                "quote": "Reliability is the foundation of autonomy.",
                "attribution": "System Architect",
            },
            {
                "type": "table",
                "title": "Regional Performance",
                "columns": ["Region", "Growth", "Status"],
                "rows": [["Americas", "+24%", "On Track"], ["EMEA", "+18%", "Ahead"]],
            },
            {
                "type": "chart",
                "title": "Quarterly Expansion",
                "chart": {
                    "kind": "bar",
                    "categories": ["Q1", "Q2", "Q3", "Q4"],
                    "series": [{"name": "ARR ($M)", "values": [4.2, 5.8, 8.1, 11.4]}],
                },
            },
            {"type": "blank", "speaker_notes": "Conclude with Q&A session."},
        ],
    }

    res = skill.execute(
        {"action": "render", "deck_spec": spec, "output_path": str(out_file)}
    )
    assert res["success"] is True
    assert res["action"] == "render"
    assert res["slide_count"] == 10
    assert os.path.exists(out_file)
    assert res["file_size_bytes"] > 1000

    # Inspect rendered file to ensure validity
    inspect_res = skill.execute({"action": "inspect", "input_path": str(out_file)})
    assert inspect_res["success"] is True
    assert inspect_res["slide_count"] == 10
    assert inspect_res["slides"][0]["title"] == "Corporate Briefing"
    assert inspect_res["slides"][9]["has_notes"] is True


def test_render_path_traversal_rejected(skill, tmp_path):
    spec = {"title": "Safe Deck", "slides": [{"type": "blank"}]}
    unsafe_path = str(tmp_path / ".." / ".." / "escape.pptx")
    res = skill.execute(
        {"action": "render", "deck_spec": spec, "output_path": unsafe_path}
    )
    assert res["success"] is False
    assert res["error_code"] == "OUTPUT_PATH_UNSAFE"

    bad_ext_path = str(tmp_path / "presentation.pdf")
    res_bad_ext = skill.execute(
        {"action": "render", "deck_spec": spec, "output_path": bad_ext_path}
    )
    assert res_bad_ext["success"] is False
    assert res_bad_ext["error_code"] == "OUTPUT_PATH_UNSAFE"


def test_inspect_nonexistent_file(skill, tmp_path):
    missing = str(tmp_path / "does_not_exist.pptx")
    res = skill.execute({"action": "inspect", "input_path": missing})
    assert res["success"] is False
    assert res["error_code"] == "INSPECT_FAILED"


def test_12_slide_investor_deck_spec(skill, tmp_path):
    """
    Acceptance criteria from Issue #276:
    An agent can validate a 12-slide investor deck spec, render a .pptx using
    the pitch template, receive warnings for one missing optional image, open
    the file with editable text and notes, and pass tests offline in CI.
    """
    out_file = tmp_path / "series_a_pitch.pptx"
    investor_deck_spec = {
        "title": "Skillware Series A Deck",
        "template_id": "pitch_v1",
        "theme": {
            "accent_color": "#6E57E0",
            "font_heading": "Calibri",
            "font_body": "Calibri",
        },
        "metadata": {
            "author": "ARPA Hellenic Logical Systems",
            "subject": "Investor Overview",
        },
        "slides": [
            {
                "type": "title",
                "title": "Skillware",
                "subtitle": "Deterministic AI Skills for Agent Runtimes",
            },
            {
                "type": "bullets",
                "title": "The Problem",
                "bullets": [
                    "Agents hallucinate tool definitions",
                    "Fragile JSON parsing breaks agent loops",
                    "Zero audit trail on sensitive actions",
                ],
            },
            {
                "type": "bullets",
                "title": "The Solution",
                "bullets": [
                    "Deterministic, offline-first skill bundles",
                    "Strict schema validation and constitution enforcement",
                    "Universal provider adapters",
                ],
            },
            {
                "type": "section",
                "title": "Market Opportunity",
                "subtitle": "Autonomous AI Agent Infrastructure",
            },
            {
                "type": "chart",
                "title": "Agent Market Growth",
                "chart": {
                    "kind": "line",
                    "categories": ["2024", "2025", "2026", "2027"],
                    "series": [
                        {"name": "Market Size ($B)", "values": [5.1, 12.8, 28.5, 52.0]}
                    ],
                },
            },
            {
                "type": "two_column",
                "title": "Competitive Advantage",
                "left": ["Monolithic Frameworks", "High latency", "Vendor lock-in"],
                "right": [
                    "Skillware",
                    "Sub-millisecond local execute",
                    "Universal across Claude, OpenAI, Gemini",
                ],
            },
            {
                "type": "table",
                "title": "Traction & Milestones",
                "columns": ["Quarter", "Skills Shipped", "Total Downloads"],
                "rows": [
                    ["Q1 2026", "12", "45,000"],
                    ["Q2 2026", "18", "120,000"],
                    ["Q3 2026", "25", "310,000"],
                ],
            },
            {
                "type": "image",
                "title": "Architecture Diagram",
                "image": {"path": "/tmp/optional_arch_diagram_missing.png"},
                "caption": "Figure: Host-to-skill boundary",
            },
            {
                "type": "quote",
                "quote": "Deterministic skills are the fundamental building blocks of production agentic software.",
                "attribution": "Lead AI Researcher",
            },
            {
                "type": "bullets",
                "title": "Business Model",
                "bullets": [
                    "Open-source core registry",
                    "Enterprise SLA & customized skills support",
                    "Private corporate registry hosting",
                ],
            },
            {
                "type": "bullets",
                "title": "The Ask",
                "bullets": [
                    "$12M Series A financing",
                    "18 months runway",
                    "Key engineering & developer advocacy hires",
                ],
                "speaker_notes": "Emphasize capital efficiency and current organic developer pull.",
            },
            {
                "type": "blank",
                "speaker_notes": "Thank the investors and open for partner questions.",
            },
        ],
    }

    # 1. Validation check
    val_res = skill.execute(
        {"action": "validate_spec", "deck_spec": investor_deck_spec}
    )
    assert val_res["valid"] is True
    assert val_res["slide_count"] == 12
    # Receives warning for the one missing optional image
    assert any(w["code"] == "ASSET_NOT_FOUND" for w in val_res["warnings"])

    # 2. Render check
    render_res = skill.execute(
        {
            "action": "render",
            "deck_spec": investor_deck_spec,
            "output_path": str(out_file),
        }
    )
    assert render_res["success"] is True
    assert render_res["slide_count"] == 12
    assert os.path.exists(out_file)

    # 3. Inspect check
    inspect_res = skill.execute({"action": "inspect", "input_path": str(out_file)})
    assert inspect_res["success"] is True
    assert inspect_res["slide_count"] == 12
    assert inspect_res["slides"][10]["has_notes"] is True
    assert inspect_res["slides"][11]["has_notes"] is True


def test_constitution_offline_no_remote_apis():
    root = os.path.dirname(__file__)
    banned = (
        "openai",
        "anthropic",
        "gemini",
        "requests.get",
        "requests.post",
        "urllib.request",
    )
    for name in (
        "skill.py",
        "builder.py",
        "placeholders.py",
        "archetypes.py",
        "lint.py",
    ):
        text = open(os.path.join(root, name), encoding="utf-8").read().lower()
        for token in banned:
            assert token not in text, f"Found banned remote token '{token}' in {name}"


def test_suggest_outline_all_archetypes(skill):
    archetypes = [
        "investor_pitch",
        "technical_brief",
        "quarterly_review",
        "product_launch",
        "training_workshop",
    ]
    for arch in archetypes:
        res = skill.execute({"action": "suggest_outline", "archetype": arch})
        assert res["success"] is True
        assert res["archetype"] == arch
        assert "deck_spec" in res
        spec = res["deck_spec"]
        assert len(spec["slides"]) >= 4
        # Validate that generated skeleton conforms to schema
        val = skill.execute({"action": "validate_spec", "deck_spec": spec})
        assert val["valid"] is True


def test_suggest_outline_negative_constraints(skill):
    # Test filtering out pricing/financials
    res_no_pricing = skill.execute(
        {
            "action": "suggest_outline",
            "archetype": "investor_pitch",
            "constraints": ["no pricing", "no business model"],
        }
    )
    assert res_no_pricing["success"] is True
    titles = [s.get("title", "").lower() for s in res_no_pricing["deck_spec"]["slides"]]
    assert not any("business model" in t for t in titles)


def test_lint_deck_perfect_score(skill):
    clean_spec = {
        "title": "Clean Presentation",
        "slides": [
            {"type": "title", "title": "Welcome", "subtitle": "Intro"},
            {
                "type": "bullets",
                "title": "Overview",
                "bullets": ["Item one", "Item two"],
            },
            {
                "type": "metrics",
                "title": "KPIs",
                "metrics": [{"value": "99.9%", "label": "Uptime"}],
            },
            {
                "type": "timeline",
                "title": "Roadmap",
                "items": [{"date": "Q1", "title": "Alpha", "description": "Release"}],
            },
        ],
    }
    res = skill.execute({"action": "lint_deck", "deck_spec": clean_spec})
    assert res["success"] is True
    assert res["score"] == 100
    assert res["passed"] is True
    assert len(res["errors"]) == 0
    assert len(res["warnings"]) == 0


def test_lint_deck_detects_violations(skill):
    dirty_spec = {
        "title": "Flawed Presentation",
        "slides": [
            {"type": "title", "title": ""},  # EMPTY_TITLE
            {
                "type": "bullets",
                "title": "Text Heavy",
                "bullets": [
                    "A" * 220,  # WALL_OF_TEXT
                ],  # ORPHAN_BULLET (only 1 item)
            },
            {
                "type": "metrics",
                "title": "Missing Labels",
                "metrics": [{"value": "$10M"}],  # METRIC_WITHOUT_LABEL
            },
        ],
    }
    res = skill.execute({"action": "lint_deck", "deck_spec": dirty_spec})
    assert res["success"] is True
    assert res["score"] < 100
    issue_codes = [i["code"] for i in res["issues"]]
    assert "EMPTY_TITLE" in issue_codes
    assert "WALL_OF_TEXT" in issue_codes
    assert "ORPHAN_BULLET" in issue_codes
    assert "METRIC_WITHOUT_LABEL" in issue_codes


def test_lint_deck_min_score_and_a11y(skill):
    spec = {
        "title": "Short Deck",
        "slides": [
            {"type": "title", "title": "Deck Title"},
            {
                "type": "image",
                "title": "Photo",
                "image": {"placeholder_id": "hero"},
            },  # MISSING_ALT in strict_a11y
        ],  # LOW_SLIDE_COUNT (<3)
    }
    # strict_a11y check
    res = skill.execute(
        {"action": "lint_deck", "deck_spec": spec, "strict_a11y": True, "min_score": 95}
    )
    assert res["success"] is True
    assert res["passed"] is False  # score should be < 95
    issue_codes = [i["code"] for i in res["issues"]]
    assert "LOW_SLIDE_COUNT" in issue_codes
    assert "MISSING_ALT" in issue_codes


def test_render_v020_new_layouts(skill, tmp_path):
    out_file = tmp_path / "v020_layouts.pptx"
    spec = {
        "title": "v0.2.0 Layout Showcase",
        "metadata": {
            "classification": "CONFIDENTIAL",
            "legal_footer": "Internal Confidential - Do Not Distribute",
        },
        "slides": [
            {
                "type": "title",
                "title": "v0.2.0 Showcase",
                "subtitle": "New Layouts and Fit Policies",
                "image": {"placeholder_id": "logo", "fit": "contain"},
            },
            {
                "type": "timeline",
                "title": "Product Roadmap",
                "items": [
                    {
                        "date": "Q1 2026",
                        "title": "Phase 1: Alpha",
                        "description": "Core layout architecture",
                        "status": "completed",
                    },
                    {
                        "date": "Q2 2026",
                        "title": "Phase 2: Beta",
                        "description": "Smart placeholders and linting",
                        "status": "in_progress",
                    },
                    {
                        "date": "Q3 2026",
                        "title": "Phase 3: GA",
                        "description": "Enterprise rollout",
                        "status": "planned",
                    },
                ],
                "speaker_notes": "Highlight phase 2 progress.",
            },
            {
                "type": "metrics",
                "title": "Quarterly Performance",
                "metrics": [
                    {
                        "value": "99.98%",
                        "label": "System Availability",
                        "delta": "+0.4%",
                        "trend": "up",
                    },
                    {
                        "value": "$4.8M",
                        "label": "Annualized Run Rate",
                        "delta": "+28% YoY",
                        "trend": "up",
                    },
                    {
                        "value": "12 ms",
                        "label": "p99 Execution Latency",
                        "delta": "-5 ms",
                        "trend": "down",
                    },
                    {
                        "value": "1,420",
                        "label": "Active Organizations",
                        "trend": "neutral",
                    },
                ],
                "speaker_notes": "Notice all key metrics beating targets.",
            },
            {
                "type": "comparison",
                "title": "Architectural Comparison",
                "left": {
                    "title": "Legacy Approach",
                    "items": [
                        "Manual slide deck authoring",
                        "Unstructured copy-paste errors",
                        "No quality linting",
                    ],
                },
                "right": {
                    "title": "Skillware Platform",
                    "items": [
                        "Deterministic code-driven generation",
                        "Automated pre-flight schema checks",
                        "Zero remote network calls",
                    ],
                },
                "speaker_notes": "Clear differentiation on deterministic reliability.",
            },
            {
                "type": "image",
                "title": "Hero Placeholder",
                "image": {
                    "placeholder_id": "hero",
                    "fit": "cover",
                    "alt": "Hero architecture diagram",
                },
                "caption": "Figure: Architectural schematic",
            },
        ],
    }

    # Validate
    val = skill.execute({"action": "validate_spec", "deck_spec": spec})
    assert val["valid"] is True

    # Lint
    lint = skill.execute({"action": "lint_deck", "deck_spec": spec})
    assert lint["passed"] is True
    assert lint["score"] >= 95

    # Render
    render = skill.execute(
        {"action": "render", "deck_spec": spec, "output_path": str(out_file)}
    )
    assert render["success"] is True
    assert render["slide_count"] == 5
    assert render["slides"][1]["type"] == "timeline"
    assert render["slides"][2]["type"] == "metrics"
    assert render["slides"][3]["type"] == "comparison"
    assert out_file.is_file()
    assert out_file.stat().st_size > 1000

    # Inspect
    inspect_res = skill.execute({"action": "inspect", "input_path": str(out_file)})
    assert inspect_res["success"] is True
    assert inspect_res["slide_count"] == 5
    assert inspect_res["slides"][1]["layout_name"] is not None
    assert inspect_res["slides"][1]["has_notes"] is True


def test_governance_ribbon_and_footer_rendered(skill, tmp_path):
    import pptx

    out_file = tmp_path / "gov_rendered.pptx"
    spec = {
        "title": "Governance Presentation",
        "metadata": {
            "classification": "RESTRICTED",
            "legal_footer": "Proprietary and Confidential - Skillware",
        },
        "slides": [
            {
                "type": "title",
                "title": "Cover Slide",
                "subtitle": "Confidential Overview",
            },
            {"type": "bullets", "title": "Key Items", "bullets": ["Item A", "Item B"]},
        ],
    }

    render = skill.execute(
        {"action": "render", "deck_spec": spec, "output_path": str(out_file)}
    )
    assert render["success"] is True
    assert out_file.is_file()

    # Inspect rendered pptx shapes directly to guarantee ribbon and footer are not no-ops
    prs = pptx.Presentation(str(out_file))
    assert len(prs.slides) == 2
    for s in prs.slides:
        texts = []
        for shape in s.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    texts.append(p.text)
        assert any(
            "[RESTRICTED]" in t for t in texts
        ), f"Missing classification ribbon in {texts}"
        assert any(
            "Proprietary and Confidential - Skillware" in t for t in texts
        ), f"Missing legal footer in {texts}"

    # Also test backward-compatible root fallback
    out_file_root = tmp_path / "gov_root_fallback.pptx"
    spec_root = {
        "title": "Root Fallback Governance",
        "classification": "INTERNAL",
        "legal_footer": "Internal Only",
        "slides": [
            {"type": "title", "title": "Root Cover"},
        ],
    }
    render_root = skill.execute(
        {"action": "render", "deck_spec": spec_root, "output_path": str(out_file_root)}
    )
    assert render_root["success"] is True
    prs_root = pptx.Presentation(str(out_file_root))
    root_texts = [
        p.text
        for s in prs_root.slides
        for shape in s.shapes
        if shape.has_text_frame
        for p in shape.text_frame.paragraphs
    ]
    assert any("[INTERNAL]" in t for t in root_texts)
    assert any("Internal Only" in t for t in root_texts)
